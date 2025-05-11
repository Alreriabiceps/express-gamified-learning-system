import sys
import json
from pptx import Presentation

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

if __name__ == "__main__":
    pptx_path = sys.argv[1]
    text = extract_text_from_pptx(pptx_path)
    print(json.dumps({"text": text})) 